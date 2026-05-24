import os
import pandas as pd
from pptx import Presentation


def export_ppt(contributions, timeline):
    os.makedirs("data", exist_ok=True)

    # 누적 CSV 읽기
    df = pd.read_csv("data/commit_report.csv")

    prs = Presentation()

    # 1. 제목 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "개발 과정 요약"
    slide.placeholders[1].text = "GitHub Commit 기반 자동 생성"

    # 2. 팀원 기여도
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "팀원 기여도"

    text = ""
    for author, count in contributions.items():
        text += f"{author}: {count} commits\n"

    slide.placeholders[1].text = text

    # 3. 개발 타임라인
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "개발 타임라인"

    text = ""
    for date, items in list(timeline.items())[:10]:
        text += f"{date}\n"
        text += f"- {items[0]}\n\n"

    slide.placeholders[1].text = text

    # 4. 커밋별 슬라이드 (전체!)
    for _, commit in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = commit["category"]

        if commit.get("is_merge", False):
            title += " (PR Merge)"

        slide.shapes.title.text = title

        body = (
            f"날짜: {commit['date']}\n"
            f"작성자: {commit['author']}\n"
            f"요약: {commit['summary']}\n"
            f"SHA: {commit['sha']}"
        )

        slide.placeholders[1].text = body

    prs.save("data/commit_report.pptx")